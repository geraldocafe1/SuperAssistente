"""
=================================================================================
SUPER ASSISTENTE MULTIMODAL v2.2 - VERSÃO REVISADA E ESTÁVEL
=================================================================================
Descrição:
Aplicação web interativa que funciona como um assistente multimodal. Esta versão
inclui uma revisão na lógica da interface para garantir que todos os uploaders
de mídia sejam exibidos corretamente.

Pré-requisitos:
pip install streamlit google-generativeai python-dotenv gtts streamlit-audiorec Pillow python-docx python-pptx pandas PyPDF2
"""

# --- 1. IMPORTAÇÃO DAS BIBLIOTECAS ---
import streamlit as st
import google.generativeai as genai
from PIL import Image
import os
from dotenv import load_dotenv
from gtts import gTTS
import io
import speech_recognition as sr
from st_audiorec import st_audiorec
import docx
from pptx import Presentation
import pandas as pd
import PyPDF2

# --- 2. CONFIGURAÇÃO INICIAL DA PÁGINA E API ---
st.set_page_config(
    page_title="Super Assistente Multimodal",
    page_icon="🤖",
    layout="wide"
)
load_dotenv()
api_key = os.getenv("GEMINI_API_KEY")
if not api_key:
    st.error("Chave de API do Gemini não encontrada! Configure seu arquivo .env.")
    st.stop()
genai.configure(api_key=api_key)

# --- 3. FUNÇÕES AUXILIARES ---
@st.cache_resource
def carregar_modelo():
    try:
        return genai.GenerativeModel('gemini-1.5-flash-latest')
    except Exception as e:
        st.error(f"Erro ao carregar o modelo: {e}")
        st.stop()

def extrair_texto_docx(arquivo):
    try:
        return "\n".join([p.text for p in docx.Document(arquivo).paragraphs])
    except Exception as e: return f"Erro ao ler .docx: {e}"

def extrair_texto_pptx(arquivo):
    try:
        return "\n".join([s.text for slide in Presentation(arquivo).slides for s in slide.shapes if hasattr(s, "text")])
    except Exception as e: return f"Erro ao ler .pptx: {e}"

def extrair_texto_csv(arquivo):
    try:
        return pd.read_csv(arquivo).to_markdown(index=False)
    except Exception as e: return f"Erro ao ler .csv: {e}"

def extrair_texto_pdf(arquivo):
    try:
        leitor = PyPDF2.PdfReader(arquivo)
        return "\n".join([pagina.extract_text() for pagina in leitor.pages])
    except Exception as e: return f"Erro ao ler .pdf: {e}"

def audio_para_texto(audio_bytes):
    r = sr.Recognizer()
    try:
        with sr.AudioFile(io.BytesIO(audio_bytes)) as source:
            audio = r.record(source)
        return r.recognize_google(audio, language='pt-BR')
    except (sr.UnknownValueError, sr.RequestError): return None

def texto_para_audio(texto):
    try:
        tts = gTTS(text=texto, lang='pt', slow=False)
        fp = io.BytesIO()
        tts.write_to_fp(fp)
        fp.seek(0)
        return fp.read()
    except Exception: return None

# --- 4. INICIALIZAÇÃO DO ESTADO DA SESSÃO ---
if "modelo" not in st.session_state:
    st.session_state.modelo = carregar_modelo()
if "historico_chat" not in st.session_state:
    st.session_state.historico_chat = []
if "contexto_ativo" not in st.session_state:
    st.session_state.contexto_ativo = None

# --- 5. LAYOUT DA INTERFACE (SIDEBAR) ---
with st.sidebar:
    st.header("Carregar Mídia para Análise")
    st.markdown("---")
    
    # Adicionei subtítulos para organizar melhor e garantir que tudo apareça.
    st.subheader("📷 Análise de Imagem")
    imagem_carregada = st.file_uploader("Clique para enviar uma Imagem", type=["jpg", "jpeg", "png"], key="uploader_imagem")
    
    st.markdown("---")
    st.subheader("🎬 Análise de Vídeo")
    video_carregado = st.file_uploader("Clique para enviar um Vídeo (curto)", type=["mp4", "mov", "avi", "mpeg"], key="uploader_video")

    st.markdown("---")
    st.subheader("📄 Análise de Documento")
    documento_carregado = st.file_uploader("Clique para enviar um Documento", type=["txt", "md", "csv", "py", "docx", "pptx", "pdf"], key="uploader_documento")
    
    st.markdown("---")
    if st.button("Limpar Mídia Carregada"):
        st.session_state.contexto_ativo = None
        st.rerun()

# --- Lógica de processamento de upload (fora da renderização da sidebar) ---
if imagem_carregada:
    st.session_state.contexto_ativo = {"tipo": "imagem", "conteudo": Image.open(imagem_carregada)}
    with st.sidebar:
        st.image(imagem_carregada, caption="Imagem Carregada")
        st.success("Imagem pronta!")
elif video_carregado:
    st.session_state.contexto_ativo = {"tipo": "vídeo", "conteudo": None} # Marca que um vídeo foi carregado
    with st.sidebar:
        st.video(video_carregado)
        with st.spinner("Processando vídeo..."):
            st.session_state.contexto_ativo["conteudo"] = genai.upload_file(path=video_carregado.name, resource=video_carregado)
        st.success("Vídeo pronto!")
elif documento_carregado:
    st.session_state.contexto_ativo = {"tipo": "documento", "conteudo": None}
    with st.sidebar, st.spinner("Lendo documento..."):
        nome_arquivo = documento_carregado.name
        if nome_arquivo.endswith(".pdf"): texto = extrair_texto_pdf(documento_carregado)
        elif nome_arquivo.endswith(".docx"): texto = extrair_texto_docx(documento_carregado)
        elif nome_arquivo.endswith(".pptx"): texto = extrair_texto_pptx(documento_carregado)
        elif nome_arquivo.endswith(".csv"): texto = extrair_texto_csv(documento_carregado)
        else: texto = documento_carregado.getvalue().decode("utf-8")
        st.session_state.contexto_ativo["conteudo"] = texto
        st.text_area("Conteúdo Extraído", texto, height=150, disabled=True)
        st.success("Documento pronto!")

# --- 6. INTERFACE PRINCIPAL DO CHAT ---
st.title("🤖 Super Assistente Multimodal")
st.caption(f"Modelo IA em uso: {st.session_state.modelo.model_name}")

for role, text in st.session_state.historico_chat:
    with st.chat_message(role): st.markdown(text)

st.write("---")
st.subheader("🎤 Interaja por Voz")
wav_audio_data = st_audiorec()

prompt = None
if wav_audio_data: prompt = audio_para_texto(wav_audio_data)
if input_text := st.chat_input("Ou digite sua mensagem..."): prompt = input_text

# --- 7. PROCESSAMENTO E RESPOSTA ---
if prompt:
    st.session_state.historico_chat.append(("user", prompt))
    with st.chat_message("user"): st.markdown(prompt)

    conteudo_para_gemini = []
    info_contexto = ""
    if st.session_state.contexto_ativo and st.session_state.contexto_ativo["conteudo"]:
        contexto = st.session_state.contexto_ativo
        tipo_contexto = contexto["tipo"]
        
        if tipo_contexto == "documento":
            prompt_formatado = f"Com base no documento fornecido, responda à pergunta.\n\n--- DOCUMENTO ---\n{contexto['conteudo']}\n--- FIM ---\n\nPergunta: {prompt}"
            conteudo_para_gemini.append(prompt_formatado)
        else:
            conteudo_para_gemini = [prompt, contexto["conteudo"]]
        info_contexto = f" (com base no {tipo_contexto} carregado)"
        st.session_state.contexto_ativo = None
    else:
        conteudo_para_gemini.append(prompt)

    with st.chat_message("model"), st.spinner(f"Gemini está pensando{info_contexto}..."):
        try:
            resposta_modelo = st.session_state.modelo.generate_content(conteudo_para_gemini)
            resposta_texto = resposta_modelo.text
        except Exception as e:
            resposta_texto = f"Desculpe, ocorreu um erro: {e}"
            st.error(resposta_texto)
    
        st.markdown(resposta_texto)
        audio_resposta = texto_para_audio(resposta_texto)
        if audio_resposta: st.audio(audio_resposta, autoplay=True)
        st.session_state.historico_chat.append(("model", resposta_texto))